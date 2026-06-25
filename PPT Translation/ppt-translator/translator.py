"""
translator.py — PPT 翻译核心引擎

流程：
  1. 用 python-pptx 解析 .pptx，逐 slide 收集所有文本段落
  2. 将每张 slide 的文本批量发给 LLM（JSON in / JSON out）
  3. 将译文逐段写回，保留原始字体/颜色/大小等格式
  4. 保存为新 .pptx 文件
"""

from __future__ import annotations

import json
import re
import time
import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, List, Optional

import httpx
from openai import OpenAI
from pptx import Presentation

from config import API_KEY, BASE_URL, MODEL, BATCH_SIZE, PROXY

logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────
# 翻译场景提示词
# ──────────────────────────────────────────────
DOMAIN_HINTS: dict[str, str] = {
    "general":  "通用内容，自然流畅即可",
    "business": "商务场景，使用正式商务表达和专业商业术语",
    "medical":  "医疗健康场景，使用精准的医学专业术语",
    "legal":    "法律合同场景，使用严谨规范的法律用语",
    "finance":  "金融投资场景，使用规范的财务和金融术语",
    "academic": "学术论文场景，符合学术写作规范",
    "tech":     "科技/IT场景，使用准确的技术术语",
}

LANG_NAMES: dict[str, str] = {
    "zh": "中文",
    "en": "英文",
}


# ──────────────────────────────────────────────
# 数据结构
# ──────────────────────────────────────────────
@dataclass
class TextUnit:
    """表示一个需要翻译的段落单元"""
    text_frame: object   # pptx text_frame 对象
    para_idx: int        # 该段落在 text_frame.paragraphs 中的索引
    original_text: str   # 原始文本


# ──────────────────────────────────────────────
# 文本收集
# ──────────────────────────────────────────────
def _collect_from_text_frame(tf, units: List[TextUnit]) -> None:
    """从 text_frame 收集所有非空段落"""
    for i, para in enumerate(tf.paragraphs):
        text = para.text.strip()
        if text:
            units.append(TextUnit(tf, i, text))


def collect_text_units(slide) -> List[TextUnit]:
    """收集幻灯片内所有可翻译的文本单元（跳过图片/图表/备注）"""
    units: List[TextUnit] = []
    for shape in slide.shapes:
        # 普通文本框
        if shape.has_text_frame:
            _collect_from_text_frame(shape.text_frame, units)
        # 表格
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _collect_from_text_frame(cell.text_frame, units)
    return units


# ──────────────────────────────────────────────
# 译文写回
# ──────────────────────────────────────────────
def write_translation(unit: TextUnit, new_text: str) -> None:
    """
    将译文写回段落，策略：
      - 把所有译文放入第一个 run（保留其字体/颜色/大小格式）
      - 其余 run 清空
    """
    para = unit.text_frame.paragraphs[unit.para_idx]
    if not para.runs:
        return  # 无 run 的段落跳过（极少见边缘情况）

    para.runs[0].text = new_text
    for run in para.runs[1:]:
        run.text = ""


# ──────────────────────────────────────────────
# LLM 翻译
# ──────────────────────────────────────────────
def _build_prompt(texts_dict: dict, source_lang: str, target_lang: str, domain: str) -> str:
    src = LANG_NAMES.get(source_lang, source_lang)
    tgt = LANG_NAMES.get(target_lang, target_lang)
    hint = DOMAIN_HINTS.get(domain, "通用内容")

    return f"""你是专业翻译。请将下方 JSON 中的{src}文本翻译为{tgt}。
翻译场景：{hint}

规则：
1. 只翻译 value，不修改 key
2. 纯数字、代码变量、URL 原样保留
3. 保持原文语气，不添加解释或注释
4. 直接返回 JSON，不要输出其他任何内容

输入：
{json.dumps(texts_dict, ensure_ascii=False)}"""


def translate_batch(
    texts: List[str],
    source_lang: str,
    target_lang: str,
    domain: str,
    client: OpenAI,
) -> List[str]:
    """
    批量翻译文本列表，返回等长的译文列表。
    遇到 LLM 返回异常时最多重试 2 次，仍失败则原文返回。
    """
    # 只翻译非空项，建立索引映射
    indexed = {str(i): t for i, t in enumerate(texts) if t.strip()}
    if not indexed:
        return texts

    prompt = _build_prompt(indexed, source_lang, target_lang, domain)

    for attempt in range(3):
        try:
            resp = client.chat.completions.create(
                model=MODEL,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.1,
            )
            raw = resp.choices[0].message.content.strip()

            # 去掉 LLM 可能包裹的 markdown 代码块
            raw = re.sub(r"^```(?:json)?\s*", "", raw)
            raw = re.sub(r"\s*```$", "", raw)

            translated: dict = json.loads(raw.strip())

            result = list(texts)
            for k, v in translated.items():
                idx = int(k)
                if 0 <= idx < len(result):
                    result[idx] = str(v)
            return result

        except (json.JSONDecodeError, ValueError, KeyError) as e:
            logger.warning(f"翻译解析失败（第 {attempt + 1} 次）：{e}")
            if attempt < 2:
                time.sleep(1 * (attempt + 1))

        except Exception as e:
            logger.error(f"LLM 调用出错：{e}")
            if attempt < 2:
                time.sleep(2)

    # 全部失败：原文返回
    logger.error("翻译批次全部失败，返回原文")
    return texts


# ──────────────────────────────────────────────
# 文件名翻译
# ──────────────────────────────────────────────
def translate_filename(stem: str, source_lang: str, target_lang: str, client: OpenAI) -> str:
    """
    翻译文件名（不含扩展名）。
    - 译为英文时：单词首字母大写，空格转下划线
    - 译为中文时：直接返回中文
    - 失败时静默回退到原文件名
    """
    # 将下划线/连字符还原为空格，便于翻译
    clean = stem.replace("_", " ").replace("-", " ")

    src = LANG_NAMES.get(source_lang, source_lang)
    tgt = LANG_NAMES.get(target_lang, target_lang)

    prompt = (
        f"将以下文件名从{src}翻译为{tgt}。"
        f"{'翻译为英文时，每个单词首字母大写，单词之间用下划线连接。' if target_lang == 'en' else ''}"
        f"只返回译文，不要加任何解释或标点。\n文件名：{clean}"
    )

    try:
        resp = client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1,
        )
        translated = resp.choices[0].message.content.strip()
        # 去掉文件系统不允许的字符
        translated = re.sub(r'[<>:"/\\|?*\n]', "", translated).strip()
        if target_lang == "en":
            # 空格统一转下划线
            translated = translated.replace(" ", "_")
        return translated if translated else stem
    except Exception as e:
        logger.warning(f"文件名翻译失败，使用原名：{e}")
        return stem


# ──────────────────────────────────────────────
# 主入口
# ──────────────────────────────────────────────
def translate_pptx(
    input_path: str,
    output_path: str,
    source_lang: str,
    target_lang: str,
    domain: str = "general",
    progress_callback: Optional[Callable[[float], None]] = None,
) -> str:
    """
    翻译整个 PPTX 文件并保存到 output_path。

    Args:
        input_path:        输入 .pptx 路径
        output_path:       输出 .pptx 路径
        source_lang:       源语言 "zh" 或 "en"
        target_lang:       目标语言 "zh" 或 "en"
        domain:            翻译场景
        progress_callback: 接收 0.0~1.0 进度的回调函数

    Returns:
        译文文件名 stem（不含扩展名），用于下载时命名输出文件
    """
    # 构建 httpx 客户端：有代理配置时使用代理，否则直连
    http_client = (
        httpx.Client(proxy=PROXY, timeout=60.0)
        if PROXY
        else httpx.Client(timeout=60.0)
    )
    client = OpenAI(api_key=API_KEY, base_url=BASE_URL, http_client=http_client)
    prs = Presentation(input_path)
    total = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides):
        units = collect_text_units(slide)

        # 按 BATCH_SIZE 分块，避免超出 context window
        for chunk_start in range(0, max(len(units), 1), BATCH_SIZE):
            chunk = units[chunk_start: chunk_start + BATCH_SIZE]
            if not chunk:
                continue

            original_texts = [u.original_text for u in chunk]
            translated_texts = translate_batch(
                original_texts, source_lang, target_lang, domain, client
            )

            for unit, new_text in zip(chunk, translated_texts):
                if new_text and new_text.strip():
                    write_translation(unit, new_text)

        if progress_callback:
            progress_callback((slide_idx + 1) / total)

    prs.save(output_path)

    # 翻译文件名
    original_stem = Path(input_path).stem.replace("_in", "")  # 去掉我们加的 _in 后缀
    translated_stem = translate_filename(original_stem, source_lang, target_lang, client)
    logger.info(f"翻译完成：{output_path}，输出文件名：{translated_stem}")
    return translated_stem
